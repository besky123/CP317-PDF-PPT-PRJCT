import tkinter as tk 
from tkinter import messagebox 
from tkinter.filedialog import askopenfile, asksaveasfilename
import collections
import collections.abc
from pptx import Presentation
import PyPDF2


root = tk.Tk()
root.title("Upload PDF")

canvas = tk.Canvas(root, width= 600, height = 300)
canvas.grid(columnspan=3)

#icon 

icon = tk.PhotoImage(file = r"C:\Users\besky\OneDrive\Desktop\CP317 Project\images\pdf_icon.png")
icon_label = tk.Label(image = icon)
icon_label.grid(column=1, row = 0)

#message 

message = tk.Label(root, text = "Click the 'Browse' Button to search for the PDF")
message.grid(columnspan=3,column=0,row=1)

# A Function used to read pdf files. 
def extractor(file_path):
    try:
        #open pdf file
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file) #assign a variable to the file handler
            num_pages = len(pdf_reader.pages) #find the number of pages in the pdf

            # iterate through each page in the pdf and assign it to a "text" variable. 
            text = ""
            for page_num in range(num_pages): 
                page = pdf_reader.pages[page_num]
                text += page.extract_text()

            return text
    # Error handling 
    except FileNotFoundError:
        print("File not found.")
    except PyPDF2._utils.PdfStreamError:
        print("Unable to read the PDF file.")

# Writes text to PPT File 
def  write_to_ppt(text, ppt_file_path):
    ppt = Presentation()
    #slide layout (Layout 1)
    slyout = ppt.slide_layouts[1]

    slide = ppt.slides.add_slide(slyout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = "Write to PPT"

    content_shape = shapes.placeholders[1]
    content_shape.text = text

    ppt.save(ppt_file_path)
    print(f"Text successfuly written to {ppt_file_path}")

#Save PPT File 
def save_ppt(text):
    ppt_file = asksaveasfilename(defaultextension=".pptx", filetypes = [("PowerPoint File", "*.pptx")])
    if ppt_file:
        write_to_ppt(text, ppt_file)
        messagebox.showinfo("Succes",f"PPT File saved to:\n{ppt_file}")
    else:
        messagebox.showwarning("Warning, File not saved.")



#open file function
def open_file():
    pdf_file= askopenfile(parent=root,mode="rb", title="choose a file", filetypes=[("PDF Files", "*.pdf")])
    if pdf_file is not None:
        #Display path and Success Message
        message.config(text=f"File Found: {pdf_file.name}")
        messagebox.showinfo("Success", "PDF File Found")

        #call extractor if found
        pdf_content = extractor(pdf_file.name)
        print(pdf_content)
        save_ppt(pdf_content)
    else:        
        messagebox.showerror("Failed","File not found")


#Browse Button

browse_btn = tk.Button(root, text= "Browse file", command = lambda:open_file(), width="15", height="2")
browse_btn.grid(column=1, row=4)

canvas = tk.Canvas(root, width = 600, height = 250)
canvas.grid(columnspan=3)

root.mainloop()
                          



"""
def import_file(): 

def save_ppt():
"""