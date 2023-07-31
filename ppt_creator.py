from pptx import Presentation
from pptx.util import Inches, Pt
import re

class PPTCreator:
    def __init__(self, ppt_path):
        self.ppt_path = ppt_path

    def create_ppt(self, text_content, image_files):
        self.prs = Presentation('apply_theme.pptx')

        # Apply a theme
        #self.prs.apply_template('theme.potx')

        # Combine text from all pages for generating the "Agenda" slide
        combined_text = "\n".join(text_content)

        # Assuming the title is the first two lines of the document
        title_text = "\n".join(combined_text.split("\n")[:2])

        # Page 1: Title
        slide_layout = self.prs.slide_layouts[0]  # Use the 'title' layout
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        # Change the font size
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)  # Change the number to the desired size in points


        # Assuming the authors are in the third and fourth line of the document
        authors_text = "\n".join(combined_text.split("\n")[2:4])

        # Page 2: Authors
        slide_layout = self.prs.slide_layouts[1]  # Use the 'title and content' layout
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Authors"
        content = slide.placeholders[1]
        content.text = authors_text

        # Page 3: Agenda
        slide_layout = self.prs.slide_layouts[1]  # Use the 'title and content' layout
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Agenda"
        content = slide.placeholders[1]
        
        # Assuming the headers are the lines starting with a single digit followed by a space (no period)
        #headers = [line for line in combined_text.split("\n") if re.match(r"^[1-9] ", line)]
        # Assuming the headers are the lines starting with a single digit or Roman numeral followed by a space (no period)
        #headers = [line for line in combined_text.split("\n") if re.match(r"^[1-9IVXLCMivxclcm] ", line)]
        # Headers are lines starting with a single digit or Roman numeral followed by a space (no period), or starting with 'Abstract'
        headers = [line for line in combined_text.split("\n") if re.match(r"^(?:[1-9IVXLCMivxclcm] |Abstract )", line)]

        content.text = "\n".join(headers)

        # Additional slides for each section
        slide_layout = self.prs.slide_layouts[1]  # Use the 'title and content' layout
        for i in range(len(headers)):
            # Split the header at the first space to remove the numbering
            _, section_title = headers[i].split(" ", 1)
            # Find the start of the section
            start_index = combined_text.find(headers[i]) + len(headers[i])
            # Find the end of the section (start of the next section or end of the document)
            end_index = combined_text.find(headers[i + 1]) if i + 1 < len(headers) else len(combined_text)
            section_text = combined_text[start_index:end_index]
            # Extract the first two sentences from the section text
            sentences = section_text.split(". ")
            # Add a new slide with the section title
            slide = self.prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = section_title
            # Add the first two sentences as separate paragraphs to the text frame
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()  # Clear the default bullet point
            for s in sentences[:2]:
                p = text_frame.add_paragraph()
                p.text = s.strip()  # Remove leading and trailing spaces
                p.level = 0  # This is a top-level bullet point
                p.space_after = Pt(14)  # Add some space after the paragraph
            # Adjust the font size of the text frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)
        
        # Extra slide: References
        slide_layout = self.prs.slide_layouts[1]  # Use the 'title and content' layout
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "References"
        content = slide.placeholders[1]
        
        combined_text_str = "\n".join(text_content)
        references_start = combined_text_str.lower().rfind('references')
        if references_start != -1:
            references_text = combined_text_str[references_start:]
            references = re.findall(r'\[\d+\].*?(?=\[\d+\]|\Z)', references_text, re.DOTALL)
        else:
            references = []
        
        text_frame = content.text_frame
        text_frame.clear()  # Clear the default bullet point
        for ref in references[:3]:  # Only include the first 3 references
            p = text_frame.add_paragraph()
            p.text = ref.strip()  # Remove leading and trailing spaces
            p.level = 0  # This is a top-level bullet point


        
        # Extract images from the PDF
        images = image_files

        # Add a new slide for each image
        for i, image in enumerate(image_files):
            image_path = f"/tmp/image_{i}.png"
            image.save(image_path)

            slide_layout = self.prs.slide_layouts[6]  # Use the 'blank' layout
            slide = self.prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(5))

        self.prs.save(self.ppt_path)


