
from pptx import Presentation
from pptx.util import Inches

#This function will import already extracted images into the ppt doc.
def img_to_ppt(ppt_file, img_paths): 

    ppt = Presentation()
    #slide layout (Layout 1)
    slyout = ppt.slide_layouts[5]

    slide = ppt.slides.add_slide(slyout)

    #img position
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(4)

    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Write to PPT"

    for img_path in  img_paths:
        shapes.add_picture(img_path, left, top, width =width, height = height)

        top += height + Inches(0.5)

    ppt.save(ppt_file)
    print(f"Text successfuly written to {ppt_file}")

img_to_ppt('ppt3.pptx', im)